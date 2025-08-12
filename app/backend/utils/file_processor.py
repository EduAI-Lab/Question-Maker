import io
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from tika import parser
import tempfile
import os

class FileProcessor:
    @staticmethod
    def extract_text_from_pdf(file_data: bytes) -> str:
        try:
            pdf_file = io.BytesIO(file_data)
            pdf_reader = PdfReader(pdf_file)
            text = []
            for page in pdf_reader.pages:
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
                except Exception as e:
                    print(f"Error extracting text from page: {str(e)}")
                    continue
            
            if not text:
                # If no text was extracted, try Tika
                return FileProcessor.extract_text_using_tika(file_data)
            
            return "\n".join(text)
        except Exception as e:
            print(f"Error extracting text from PDF: {str(e)}")
            return FileProcessor.extract_text_using_tika(file_data)

    @staticmethod
    def extract_text_from_docx(file_data: bytes) -> str:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                temp_file.write(file_data)
                temp_file.flush()
                doc = Document(temp_file.name)
                text = []
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        text.append(paragraph.text)
                os.unlink(temp_file.name)
                
                if not text:
                    return FileProcessor.extract_text_using_tika(file_data)
                
                return "\n".join(text)
        except Exception as e:
            print(f"Error extracting text from DOCX: {str(e)}")
            return FileProcessor.extract_text_using_tika(file_data)

    @staticmethod
    def extract_text_from_pptx(file_data: bytes) -> str:
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                temp_file.write(file_data)
                temp_file.flush()
                prs = Presentation(temp_file.name)
                text = []
                for slide in prs.slides:
                    slide_text = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text)
                    if slide_text:
                        text.append("\n".join(slide_text))
                os.unlink(temp_file.name)
                
                if not text:
                    return FileProcessor.extract_text_using_tika(file_data)
                
                return "\n\n".join(text)
        except Exception as e:
            print(f"Error extracting text from PPTX: {str(e)}")
            return FileProcessor.extract_text_using_tika(file_data)

    @staticmethod
    def extract_text_using_tika(file_data: bytes) -> str:
        try:
            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                temp_file.write(file_data)
                temp_file.flush()
                parsed = parser.from_file(temp_file.name)
                os.unlink(temp_file.name)
                content = parsed.get("content", "").strip()
                return content if content else ""
        except Exception as e:
            print(f"Error extracting text using Tika: {str(e)}")
            return ""

    @staticmethod
    def process_file(file_data: bytes, file_extension: str) -> str:
        processors = {
            'pdf': FileProcessor.extract_text_from_pdf,
            'docx': FileProcessor.extract_text_from_docx,
            'pptx': FileProcessor.extract_text_from_pptx,
            'txt': lambda data: data.decode('utf-8'),
        }
        
        processor = processors.get(file_extension.lower())
        if not processor:
            raise ValueError(f"Unsupported file type: {file_extension}")
        
        extracted_text = processor(file_data)
        if not extracted_text:
            raise ValueError("No text could be extracted from the file")
        
        return extracted_text